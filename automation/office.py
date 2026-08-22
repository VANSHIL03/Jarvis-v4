"""
JARVIS v4 - Office Documents & PDF Processing
Supports Microsoft Word (.docx), Excel (.xlsx), PowerPoint (.pptx), and PDF extraction/summarization.
"""

import os
from pathlib import Path
from typing import List, Dict, Any, Optional
from utils.logger import logger

class OfficeAutomation:
    def create_word_document(self, file_path: str, title: str, content_paragraphs: List[str]) -> bool:
        """Creates a Word (.docx) document with title and paragraphs."""
        try:
            import docx
            doc = docx.Document()
            doc.add_heading(title, 0)
            for p in content_paragraphs:
                doc.add_paragraph(p)
            doc.save(file_path)
            logger.info(f"Word document created: {file_path}")
            return True
        except ImportError:
            logger.warning("python-docx not installed. Creating plain text document fallback.")
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(f"# {title}\n\n" + "\n\n".join(content_paragraphs))
            return True
        except Exception as e:
            logger.error(f"Failed to create Word document: {e}")
            return False

    def create_excel_sheet(self, file_path: str, headers: List[str], rows: List[List[Any]]) -> bool:
        """Creates an Excel (.xlsx) workbook with table headers and rows."""
        try:
            import openpyxl
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.append(headers)
            for row in rows:
                ws.append(row)
            wb.save(file_path)
            logger.info(f"Excel sheet created: {file_path}")
            return True
        except ImportError:
            logger.warning("openpyxl not installed. Creating CSV fallback.")
            import csv
            with open(file_path, "w", newline="", encoding="utf-8") as f:
                writer = csv.writer(f)
                writer.writerow(headers)
                writer.writerows(rows)
            return True
        except Exception as e:
            logger.error(f"Failed to create Excel sheet: {e}")
            return False

    def create_powerpoint_presentation(self, file_path: str, slides_data: List[Dict[str, str]]) -> bool:
        """Creates a PowerPoint (.pptx) presentation from slide titles & bullets."""
        try:
            from pptx import Presentation
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[1]

            for slide_info in slides_data:
                slide = prs.slides.add_slide(blank_slide_layout)
                slide.shapes.title.text = slide_info.get("title", "Slide")
                slide.placeholders[1].text = slide_info.get("content", "")

            prs.save(file_path)
            logger.info(f"PowerPoint presentation created: {file_path}")
            return True
        except ImportError:
            logger.warning("python-pptx not installed.")
            return False
        except Exception as e:
            logger.error(f"Failed to create PowerPoint: {e}")
            return False

    def read_pdf(self, file_path: str) -> str:
        """Extracts plain text from PDF file."""
        if not os.path.exists(file_path):
            return ""

        try:
            import pypdf
            reader = pypdf.PdfReader(file_path)
            text = ""
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
            return text.strip()
        except Exception:
            try:
                import pdfplumber
                text = ""
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        text += (page.extract_text() or "") + "\n"
                return text.strip()
            except Exception as ex:
                logger.error(f"pdfplumber failed: {ex}")
                return ""

    def create_pdf(self, file_path: str, title: str = "", body: str = "") -> Dict[str, Any]:
        """
        Writes a PDF.

        reportlab is used when present. Without it the text is still saved -- as a
        .docx via python-docx, or a .txt as a last resort -- rather than failing,
        because losing a summary the user just asked for is worse than getting it
        in a different container.
        """
        target = Path(file_path)
        paragraphs = [p.strip() for p in str(body or "").split("\n") if p.strip()]

        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib.units import mm
            from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
            from xml.sax.saxutils import escape

            target.parent.mkdir(parents=True, exist_ok=True)
            styles = getSampleStyleSheet()
            story = []
            if title:
                story.append(Paragraph(escape(str(title)), styles["Title"]))
                story.append(Spacer(1, 6 * mm))
            for para in paragraphs or ["(empty document)"]:
                story.append(Paragraph(escape(para), styles["BodyText"]))
                story.append(Spacer(1, 3 * mm))

            SimpleDocTemplate(
                str(target), pagesize=A4,
                leftMargin=18 * mm, rightMargin=18 * mm,
                topMargin=18 * mm, bottomMargin=18 * mm,
            ).build(story)
            logger.info(f"PDF created: {target}")
            return {"status": "success", "path": str(target), "engine": "reportlab"}
        except ImportError:
            logger.warning("reportlab not installed; falling back to .docx / .txt output.")
        except Exception as e:
            logger.error(f"reportlab PDF creation failed: {e}")

        docx_path = target.with_suffix(".docx")
        if self.create_word_document(str(docx_path), title or target.stem, paragraphs):
            return {
                "status": "success",
                "path": str(docx_path),
                "engine": "docx-fallback",
                "message": "reportlab missing, saved as .docx instead of .pdf.",
            }

        txt_path = target.with_suffix(".txt")
        try:
            txt_path.parent.mkdir(parents=True, exist_ok=True)
            txt_path.write_text(
                (f"{title}\n\n" if title else "") + "\n\n".join(paragraphs), encoding="utf-8"
            )
            return {
                "status": "success",
                "path": str(txt_path),
                "engine": "text-fallback",
                "message": "No PDF or DOCX writer available, saved as .txt.",
            }
        except Exception as e:
            logger.error(f"Failed to write any document for '{file_path}': {e}")
            return {"status": "error", "message": str(e)}

